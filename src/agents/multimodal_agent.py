import os
from openai import OpenAI
from pathlib import Path
import streamlit as st
from datetime import datetime
from config.settings import settings

# Initialize OpenAI Client
client = OpenAI(api_key=settings.openai_api_key)

TRANSCRIPT_DIR = Path("data/generated/transcripts")
if not TRANSCRIPT_DIR.exists():
    TRANSCRIPT_DIR.mkdir(parents=True, exist_ok=True)

import base64

def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def describe_image(image_path: str) -> str:
    """
    Generates a detailed description of an image using GPT-4o Vision.
    Acts as the 'transcript' for downstream RAG tasks.
    """
    path = Path(image_path)
    cache_path = TRANSCRIPT_DIR / f"{path.name}.txt"
    
    if cache_path.exists():
        return cache_path.read_text()
        
    try:
        base64_image = encode_image(image_path)
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe this image in extreme detail. If it contains a chart, graph, or table, EXPLICITLY extract the numerical data, axis labels, and legend into a structured list or CSV format within the description. Identify key trends. Include all visible text, objects, people, emotions, colors, and layout. Ensure specific chart data is extracted if available."},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ],
                }
            ],
            max_tokens=4000
        )
        description = response.choices[0].message.content
        
        # Cache it
        cache_path.write_text(description)
        return description
        
    except Exception as e:
        return f"Error describing image: {str(e)}"

def transcribe_audio(file_path: str) -> str:
    """
    Transcribes audio/video file using OpenAI Whisper API.
    Returns the transcript text.
    Cache mechanism: Checks if {filename}.txt exists in data/generated/transcripts.
    """
    path = Path(file_path)
    if not path.exists():
        return f"Error: File not found at {file_path}"

    # Check cache
    transcript_file = TRANSCRIPT_DIR / f"{path.name}.txt"
    if transcript_file.exists():
        with open(transcript_file, "r") as f:
            return f.read()

    try:
        # Pre-process: If video, distinct audio extraction
        processed_path = path
        is_temp = False
        
        # Check size (OpenAI limit ~25MB)
        file_size_mb = path.stat().st_size / (1024 * 1024)
        
        # If video or large audio, try to extract/compress
        if path.suffix.lower() in ['.mp4', '.mov', '.avi', '.mkv'] or file_size_mb > 24:
            st.info(f"Processing media ({file_size_mb:.1f} MB)... Extracting audio for transcription.")
            try:
                # MoviePy v2.0+ compatibility
                try:
                    from moviepy import VideoFileClip
                except ImportError:
                    from moviepy.editor import VideoFileClip
                    
                video = VideoFileClip(str(path))
                if video.audio is None:
                    return "Error: This video file has no audio track to transcribe."
                    
                audio_clip = video.audio
                # Save as mp3 (compressed)
                temp_audio_path = TRANSCRIPT_DIR / f"temp_{path.stem}.mp3"
                audio_clip.write_audiofile(str(temp_audio_path), logger=None)
                
                # Close handlers
                audio_clip.close()
                video.close()
                
                processed_path = temp_audio_path
                is_temp = True
            except ImportError as e:
                return f"Error: 'moviepy' import failed. Details: {e}"
            except Exception as e:
                return f"Error extracting audio: {str(e)}"

        # Double check size of processed file
        final_size = processed_path.stat().st_size / (1024 * 1024)
        if final_size > 25:
            if is_temp: os.remove(processed_path)
            return f"Error: File is too large ({final_size:.1f} MB) even after compression. OpenAI Whisper limit is 25MB."

        with open(processed_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="whisper-1", 
                file=audio_file
            )
        
        text = transcript.text
        
        # Clean up temp
        if is_temp and processed_path.exists():
            os.remove(processed_path)
        
        # Save to cache
        with open(transcript_file, "w") as f:
            f.write(text)
            
        return text

    except Exception as e:
        return f"Error during transcription: {str(e)}"

def analyze_transcript(transcript_text: str, analysis_type: str = "summary") -> str:
    """
    Analyzes the transcript using GPT-4o to generate Summary or Minutes.
    """
    if not transcript_text or len(transcript_text) < 10:
        return "Transcript is too short to analyze."

    system_prompt = "You are an expert AI assistant capable of analyzing transcripts from meetings, calls, or videos."
    
    if analysis_type == "summary":
        user_prompt = (
            f"Please provide a concise summary of the following transcript.\n"
            f"Highlight the main topic, key arguments, and overall conclusion.\n\n"
            f"Transcript:\n{transcript_text[:15000]}" # Limit context if needed, though 4o handles 128k
        )
    elif analysis_type == "minutes":
        current_date = datetime.now().strftime("%Y-%m-%d")
        user_prompt = (
            f"Please generate structured Meeting Minutes from the following transcript.\n"
            f"Include:\n"
            f"1. 📅 Date/Time (Use '{current_date}' if not explicitly mentioned in text)\n"
            f"2. 👥 Participants (inferred)\n"
            f"3. 📝 Executive Summary\n"
            f"4. 🔑 Key Discussion Points (Bulleted)\n"
            f"5. ✅ Action Items (Who needs to do what)\n\n"
            f"Transcript:\n{transcript_text[:15000]}"
        )
    else:
        return "Unknown analysis type."

    try:
        response = client.chat.completions.create(
            model="gpt-4o", # Use high capability model for analysis
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error during analysis: {str(e)}"

def auto_translate(text: str) -> str:
    """
    Detects language and translates (En -> Fr / Fr -> En).
    """
    if not text or len(text) < 5:
        return "Text too short to translate."
        
    prompt = (
        "You are a professional translator. Task:\n"
        "1. Detect the language of the following text (English, French, or Chinese).\n"
        "2. If English, translate to French.\n"
        "3. If French, translate to English.\n"
        "4. If Chinese, translate to **BOTH English AND French**.\n"
        "5. If another language, translate to English.\n"
        "6. Output format:\n"
        "**Detected Language:** [Lang]\n\n"
        "**Translation:**\n"
        "[Translated Text (if En/Fr)]\n"
        "[**English Translation:**\nText...\n\n**French Translation:**\nText... (if Chinese)]\n\n"
        "IMPORTANT: Translate the ENTIRE text word-for-word. Do NOT summarize or truncate. Maintain all original meaning and length.\n\n"
        f"Text:\n{text[:25000]}" # Increased limit for fuller context
    )
    
    print(f"AUDIT: auto_translate called for text (first 50 chars): {text[:50]}")
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini", # 4o-mini is efficient for translation
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error during translation: {str(e)}"

    except Exception as e:
        return f"Error during translation: {str(e)}"

# --- DOCX Generation ---
from docx import Document

def save_translation_to_docx(text: str, filename: str) -> str:
    """
    Saves the translation text to a .docx file.
    Returns the path to the saved file.
    """
    if not text: return "No text to save."
    
    trans_dir = Path("data/generated/translations")
    trans_dir.mkdir(parents=True, exist_ok=True)
    save_path = trans_dir / f"{filename}_translation.docx"
    
    try:
        doc = Document()
        doc.add_heading(f"Translation: {filename}", level=1)
        doc.add_paragraph(text)
        doc.save(str(save_path))
        return str(save_path)
    except Exception as e:
        return f"Error creating DOCX: {str(e)}"

# --- Web Scraping ---
import requests
from bs4 import BeautifulSoup

def scrape_web_content(url: str) -> str:
    """
    Scrapes text content from a URL.
    Returns the text or an error message.
    """
    if not url: return "No URL provided."
    
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }
    
    try:
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Extract meaningful text (p, h1-h6, li)
        text_elements = []
        
        # Title
        if soup.title:
            text_elements.append(f"Title: {soup.title.string}\n")
            
        # Paragraphs and Headers
        for tag in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'li']):
            text = tag.get_text(strip=True)
            if len(text) > 20: # Filter short noise
                text_elements.append(text)
                
        full_text = "\n\n".join(text_elements)
        
        if len(full_text) < 100:
            return "Error: Could not extract sufficient content. The site might be JS-rendered or blocked."
            
        return full_text
        
    except requests.exceptions.RequestException as e:
        return f"Error fetching URL: {str(e)}\n\n(It might be blocked or require authentication.)"
    except Exception as e:
        return f"Error scraping: {str(e)}"

# --- RAG / Indexing Logic ---

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from agents.llm_router import init_llm

VECTOR_STORE_DIR = Path("data/vector_store")
if not VECTOR_STORE_DIR.exists():
    VECTOR_STORE_DIR.mkdir(parents=True, exist_ok=True)

def create_transcript_index(transcript_text: str, filename: str):
    """
    Creates (or overwrites) a local FAISS index for the given transcript.
    """
    if not transcript_text: return
    
    # 1. Chunking
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_text(transcript_text)
    
    if not chunks: return
    
    # 2. Embedding & Indexing
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    
    # 3. Save locally
    save_path = VECTOR_STORE_DIR / f"{filename}_index"
    vector_store.save_local(str(save_path))
    return str(save_path)

def query_transcript_index(query: str, filename: str, model_name: str = "openai:gpt-4o-mini") -> str:
    """
    Queries the specific FAISS index for the transcript and answers the question.
    """
    print(f"AUDIT: query_transcript_index called with query: {query}, file: {filename}")
    
    index_path = VECTOR_STORE_DIR / f"{filename}_index"
    if not index_path.exists():
        return "Index not found. Please transcribe the file first."
        
    try:
        # 1. Load Index
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vector_store = FAISS.load_local(str(index_path), embeddings, allow_dangerous_deserialization=True)
        
        # 2. Retrieve
        docs = vector_store.similarity_search(query, k=4)
        context = "\n\n".join([d.page_content for d in docs])
        
        # 3. Generate Answer (using selected LLM)
        llm = init_llm(model_name)
        
        system_prompt = (
            "You are a helpful analyst assistant for a Fraud Investigation platform. "
            "Use the provided context to answer the user's question directly.\n"
            "CRITICAL RULES:\n"
            "1. Answer ONLY using the provided context. If the answer is not in the context, say 'I don't find that information in the transcript.'\n"
            "2. DO NOT TRANSLATE the question. Answer the question in the SAME language it was asked.\n"
            "3. DO NOT use words like 'English Translation' or 'French Translation' unless they are in the actual content.\n"
            "4. Be concise and factual."
        )
        user_prompt = f"Context:\n{context}\n\nQuestion: {query}\nAnswer:"
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        print(f"AUDIT: Invoking LLM {model_name} for RAG")
        response = llm.invoke(messages)
        print(f"AUDIT: RAG Response received (first 50 chars): {response.content[:50]}")
        return response.content
        
    except Exception as e:
        print(f"AUDIT: ERROR in query_transcript_index: {str(e)}")
        return f"Error querying index: {str(e)}"

def evaluate_retrieval_confidence(query: str, filename: str) -> dict:
    """
    Performs a similarity search and returns retrieval scores (L2 distance) and confidence metadata.
    Lower L2 distance = detection of closer match.
    """
    index_path = VECTOR_STORE_DIR / f"{filename}_index"
    if not index_path.exists():
        return {"error": "Index not found."}
        
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vector_store = FAISS.load_local(str(index_path), embeddings, allow_dangerous_deserialization=True)
        
        # Get results with scores (lower score is better for L2 distance in FAISS)
        docs_and_scores = vector_store.similarity_search_with_score(query, k=3)
        
        results = []
        for doc, score in docs_and_scores:
            results.append({
                "content": doc.page_content[:200] + "...",
                "score": float(score) # L2 distance
            })
            
        # Heuristic confidence based on top score
        top_score = results[0]['score'] if results else 100
        # Arbitrary thresholds for "text-embedding-3-small" (L2)
        # Typically < 0.8 is good, < 0.5 is very close
        if top_score < 0.6: confidence = "High"
        elif top_score < 1.0: confidence = "Medium"
        else: confidence = "Low"
        
        return {
            "top_results": results,
            "overall_confidence": confidence,
            "best_score": top_score
        }
    except Exception as e:
        return {"error": str(e)}

# --- Aware RAG Analysis Functions ---

def analyze_speaker_aware(text: str) -> str:
    """Analyzes speakers, dialogue patterns, and key contributors."""
    if not text: return "No text."
    prompt = (
        "Perform a Speaker-Aware Analysis on the text.\n"
        "Identify:\n"
        "1. **Main Speakers**: Who is speaking? (Inferred if nameless)\n"
        "2. **Dialogue Patterns**: Who dominates? Is it a debate, interview, or monologue?\n"
        "3. **Key Contributions**: What did each speaker add?\n"
        "4. **Response**: If no dialogue exists, state that reasonable response.\n\n"
        f"Text:\n{text[:15000]}"
    )
    try:
        response = client.chat.completions.create(
            model="gpt-4o", messages=[{"role": "user", "content": prompt}], temperature=0.5
        )
        return response.choices[0].message.content
    except Exception as e: return str(e)

def analyze_time_aware(text: str) -> str:
    """Analyzes timelines, dates, and durations."""
    if not text: return "No text."
    current_date = datetime.now().strftime("%Y-%m-%d")
    prompt = (
        "Perform a Time-Aware Analysis on the text.\n"
        "Reference Date: {current_date}\n"
        "Identify:\n"
        "1. **Timeline**: Specific dates/times mentioned.\n"
        "2. **Sequence of Events**: Chronological order.\n"
        "3. **Time Range**: Duration of events discussed.\n"
        "4. **Urgency/Pacing**: Is the content urgent or historical?\n\n"
        f"Text:\n{text[:15000]}"
    )
    try:
        response = client.chat.completions.create(
            model="gpt-4o", messages=[{"role": "user", "content": prompt}], temperature=0.5
        )
        return response.choices[0].message.content
    except Exception as e: return str(e)

def analyze_confidence_aware(text: str) -> str:
    """Analyzes content quality, clarity, and reliability."""
    if not text: return "No text."
    prompt = (
        "Perform a Confidence/Quality-Aware Analysis.\n"
        "Assess:\n"
        "1. **Content Clarity**: Is the text coherent? (Score 1-10)\n"
        "2. **Information Density**: High, Medium, or Low fluff?\n"
        "3. **Ambiguity Check**: Are there vague statements?\n"
        "4. **Source Reliability**: (Inferred) Does it sound professional or informal?\n"
        "5. **Coverage quality**: Does it cover the expected topic depth?\n\n"
        f"Text:\n{text[:15000]}"
    )
    try:
        response = client.chat.completions.create(
            model="gpt-4o", messages=[{"role": "user", "content": prompt}], temperature=0.5
        )
        return response.choices[0].message.content
    except Exception as e: return str(e)

# --- PPT Generation Logic ---

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import json
from datetime import datetime
import graphviz

def analyze_sentiment(text: str) -> str:
    """
    Analyzes the text for sentiment, tone, and key insights.
    """
    if not text: return "No text provided."
    
    prompt = (
        "Analyze the following text for Sentiment and Tone.\n"
        "Provide:\n"
        "1. **Overall Sentiment**: (Positive/Negative/Neutral)\n"
        "2. **Emotional Tone**: (e.g., Formal, Angry, Excited, Concerned)\n"
        "3. **Key Emotional Drivers**: What is driving this sentiment?\n"
        "4. **AI Strategic Insight**: Hidden implications or recommended actions.\n\n"
        f"Text:\n{text[:10000]}"
    )
    
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error analyzing sentiment: {str(e)}"

def generate_workflow_diagram(text: str, filename: str) -> str:
    """
    Generates a workflow diagram (PNG) based on the logic in the text using Graphviz.
    Returns path to the generated image.
    """
    if not text: return "No text provided."
    
    # 1. Get DOT code from LLM
    prompt = (
        "You are a process mapping expert. Extract the workflow/process logic from the text below "
        "and generate Valid Graphviz DOT code to visualize it.\n"
        "Rules:\n"
        "- Use 'digraph G { ... }' format.\n"
        "- Use distinct shapes: 'box' for steps, 'diamond' for decisions, 'ellipse' for start/end.\n"
        "- Return ONLY the DOT code inside a ```dot ... ``` block or plain text.\n"
        "- Do not include markdown formatting if possible, just the code.\n\n"
        f"Text:\n{text[:10000]}"
    )
    
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2
        )
        content = response.choices[0].message.content
        
        # Cleanup markdown
        if "```dot" in content:
            dot_code = content.split("```dot")[1].split("```")[0].strip()
        elif "```" in content:
            dot_code = content.split("```")[1].split("```")[0].strip()
        else:
            dot_code = content.strip()
            
        # 2. Render with Graphviz
        img_dir = Path("data/generated/workflow")
        img_dir.mkdir(parents=True, exist_ok=True)
        
        # Create a Source object
        src = graphviz.Source(dot_code)
        output_path = img_dir / f"{filename}_workflow"
        
        # Render (returns the path including extension, usually .png)
        # format='png' ensures .png extension
        final_path = src.render(filename=str(output_path), format='png', cleanup=True)
        return final_path
        
    except Exception as e:
        return f"Error generating workflow: {str(e)}"

def generate_ppt(transcript: str, filename: str, workflow_img_path: str = None) -> tuple[str, dict]:
    """
    Generates a PowerPoint presentation from the transcript.
    Returns (path_to_pptx, content_dict).
    """
    if not transcript: return "Error: No transcript provided.", {}
    
    ppt_dir = Path("data/generated/ppt")
    ppt_dir.mkdir(parents=True, exist_ok=True)
    save_path = ppt_dir / f"{filename}.pptx"
    
    # 1. Extract Content via LLM
    try:
        system_prompt = "You are a presentation expert. Your goal is to extract structured content for a PowerPoint presentation."
        user_prompt = (
            f"Analyze the transcript and provide a JSON response with the following structure:\n"
            f"{{\n"
            f"  'title': 'Presentation Title',\n"
            f"  'subtitle': 'Subtitle (e.g. key speaker or context)',\n"
            f"  'agenda': ['Item 1', 'Item 2', 'Item 3'],\n"
            f"  'key_points': [\n"
            f"    {{'slide_title': 'Topic 1', 'points': ['Bullet A', 'Bullet B']}},\n"
            f"    {{'slide_title': 'Topic 2', 'points': ['Bullet C', 'Bullet D']}}\n"
            f"  ],\n"
            f"  'ai_insight': ['Insight 1: Hidden pattern...', 'Insight 2: Sentiment analysis...'],\n"
            f"  'conclusion': 'Final summarized takeaway',\n"
            f"  'chart_data': {{'title': 'Chart Title', 'categories': ['Cat1', 'Cat2'], 'values': [10, 20]}} \n"
            f"}}\n"
            f"IMPORTANT: If the transcript contains chart/graph data, extract strictly the real values. "
            f"If NO chart data is present, set 'chart_data' to null. "
            f"DO NOT use the example values (Cat1, 10, 20).\n\n"
            f"Transcript:\n{transcript[:15000]}"
        )
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        content = json.loads(response.choices[0].message.content)
        
    except Exception as e:
        return f"Error gathering content: {str(e)}"

    # 2. Build PPT
    try:
        prs = Presentation()
        current_date_str = datetime.now().strftime("%B %d, %Y")
        
        # Helper: Theme Application
        def apply_theme(slide, title_text, slide_num):
            # 1. Title Styling (Bold)
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    paragraph.font.bold = True
            
            # 2. Branding Bar (Top)
            width = prs.slide_width
            height = Inches(0.2)
            left = top = Inches(0)
            shape = slide.shapes.add_shape(
                1, left, top, width, height # 1 = msoShapeRectangle
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 80, 156) # Professional Blue
            shape.line.fill.background()
            
            # 3. Slide Number (Bottom Right)
            txBox = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(1), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_num)
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.RIGHT

        slide_count = 1

        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Custom Title Page Theme
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content.get('title', 'Presentation')
        for p in title.text_frame.paragraphs:
            p.font.bold = True
            
        # Add "Presented By" info
        presenter = content.get('presented_by', 'AI Assistant')
        subtitle.text = f"{content.get('subtitle', 'Generated by AI')}\n\nPresented by: {presenter}\nDate: {current_date_str}"
        
        # Add decorative bar
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 80, 156)
        
        slide_count += 1

        # Slide 2: Agenda
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        apply_theme(slide, "Agenda", slide_count)
        
        tf = slide.shapes.placeholders[1].text_frame
        for item in content.get('agenda', []):
            p = tf.add_paragraph()
            p.text = item
        slide_count += 1

        # Slides 3.x: Key Points
        for section in content.get('key_points', []):
            slide = prs.slides.add_slide(bullet_slide_layout)
            apply_theme(slide, section.get('slide_title', 'Key Metrics'), slide_count)
            
            tf = slide.shapes.placeholders[1].text_frame
            for point in section.get('points', []):
                p = tf.add_paragraph()
                p.text = point
            slide_count += 1
                
        # Slide 4: Chart
        chart_info = content.get('chart_data', {})
        if chart_info:
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
            apply_theme(slide, chart_info.get('title', 'Analysis Data'), slide_count)
            
            chart_data = CategoryChartData()
            chart_data.categories = chart_info.get('categories', ['A', 'B', 'C'])
            chart_data.add_series('Metrics', chart_info.get('values', [1, 2, 3]))
            
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            
            # Enhancement: Legend
            chart.has_legend = True
            chart.legend.include_in_layout = False
            
            # Enhancement: Axis Titles
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = chart_info.get('x_axis', 'Category')
            
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = chart_info.get('y_axis', 'Value')
            
            slide_count += 1
            
        # Slide 4.5: AI Insight Analysis
        ai_insight = content.get('ai_insight', [])
        if ai_insight:
            slide = prs.slides.add_slide(bullet_slide_layout)
            apply_theme(slide, "AI Insight Analysis", slide_count)
            tf = slide.shapes.placeholders[1].text_frame
            for point in ai_insight:
                p = tf.add_paragraph()
                p.text = point
            slide_count += 1

        # Slide: Workflow Diagram (if provided)
        if workflow_img_path and os.path.exists(workflow_img_path):
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
            apply_theme(slide, "Process Workflow", slide_count)
            
            # Add Image centered
            left = Inches(1)
            top = Inches(2)
            height = Inches(4.5)
            # Width auto-scaled or max 8 inches
            slide.shapes.add_picture(workflow_img_path, left, top, height=height)
            
            slide_count += 1
            
        # Slide 5: Conclusion
        slide = prs.slides.add_slide(bullet_slide_layout)
        apply_theme(slide, "Conclusion", slide_count)
        
        tf = slide.shapes.placeholders[1].text_frame
        p = tf.add_paragraph()
        p.text = content.get('conclusion', 'End')

        prs.save(str(save_path))
        return str(save_path), content
        
    except Exception as e:
        return f"Error creating PPT: {str(e)}", {}


def generate_audio_narration(content: dict, voice_style: str, filename: str) -> str:
    """
    Generates an audio narration (MP3) based on the PPT content.
    voice_style: 'Gentleman' (Onyx) or 'Lady' (Nova).
    """
    if not content: return "Error: No content available for narration."
    
    audio_dir = Path("data/generated/audio")
    audio_dir.mkdir(parents=True, exist_ok=True)
    save_path = audio_dir / f"{filename}_narration.mp3"
    
    # 1. Map Voice
    # OpenAI voices: alloy, echo, fable, onyx, nova, shimmer
    voice_map = {
        "Gentleman": "onyx", # Authoritative, deep male
        "Lady": "nova"       # Energetic, bright female
    }
    selected_voice = voice_map.get(voice_style, "alloy")
    
    # 2. Construct Script from Content
    # We create a natural flow
    script_parts = []
    
    # Intro
    script_parts.append(f"Welcome to the presentation titled: {content.get('title', 'Presentation')}.")
    script_parts.append(f"{content.get('subtitle', '')}.")
    
    # Agenda
    agenda = content.get('agenda', [])
    if agenda:
        script_parts.append("Today's agenda covers: " + ", ".join(agenda) + ".")
        
    # Key Points
    for section in content.get('key_points', []):
        title = section.get('slide_title', 'Topic')
        points = section.get('points', [])
        script_parts.append(f"Moving on to {title}.")
        if points:
            script_parts.append("Key takeaways include: " + ". ".join(points) + ".")
            
    # AI Insight
    insights = content.get('ai_insight', [])
    if insights:
        script_parts.append("Here are some AI-driven strategic insights.")
        script_parts.append(" ".join(insights))
        
    # Chart
    chart_info = content.get('chart_data', {})
    if chart_info:
        script_parts.append(f"Our analysis data shows meaningful trends in {chart_info.get('title', 'metrics')}.")
    
    # Conclusion
    conclusion = content.get('conclusion', '')
    if conclusion:
        script_parts.append(f"In conclusion: {conclusion}")
        
    script_parts.append("Thank you for listening.")
    
    full_script = " ".join(script_parts)
    
    # 3. Generate Audio
    try:
        response = client.audio.speech.create(
            model="tts-1",
            voice=selected_voice,
            input=full_script
        )
        response.stream_to_file(save_path)
        return str(save_path)
        
    except Exception as e:
        return f"Error generating audio: {str(e)}"
